"""Live Baizhi E2E — LLM + WebSearch MCP + pptx skill + real PPT artifact.

**双重身份**:
1. **回归测试** —— pptx skill + DashScope WebSearch + baizhi 真 LLM 整链
   端到端跑通,产出 ≥5 张幻灯片含 "Anthropic" 的真 .pptx。
2. **新 ergonomics 的展示样板** —— 用 `Agent` + `LiteLlm` + `McpToolset.http()`
   工厂,而不是手写 provider adapter + Runner / RunRequest 拼装。Stage 5
   修订(spec § 17 / § 7.5.2 / § 7.5.3)后这是 hello-world 的标准形态;原来
   ~500 行的 test 现在 ~200 行,业务量(_PptxDeckToolset)反而更显眼。

opt-in 跑法(spends real LLM / WebSearch quota):

    RUN_BAIZHI_E2E=1 .venv/bin/python -m pytest \\
      tests/test_baizhi_e2e_live_pptx_websearch.py -q -s

Keys 从 baizhi-agent/.baizhi-agent/config/.env 加载,不打印。
"""

from __future__ import annotations

import asyncio
import json
import os
from pathlib import Path
from typing import Any

import pytest

from agent_kit import Agent
from agent_kit.contrib.providers.litellm import LiteLlm
from agent_kit.contrib.skills import FilesystemSkillRegistry
from agent_kit.mcp import McpToolset
from agent_kit.provider import ToolSchema
from agent_kit.skill import SkillCatalogToolset
from agent_kit.toolset import BaseToolset, ToolCallContext
from agent_kit.types import ToolCall, ToolResult


BAIZHI_AGENT_ROOT = Path(__file__).resolve().parents[2] / "baizhi-agent"
BAIZHI_ENV_FILE = BAIZHI_AGENT_ROOT / ".baizhi-agent" / "config" / ".env"
BUNDLED_SKILLS_ROOT = BAIZHI_AGENT_ROOT / "bundled_skills"
WEBSEARCH_URL = "https://dashscope.aliyuncs.com/api/v1/mcps/WebSearch/mcp"
TASK = "深度搜索anthropic 关于AI Native orgnization组织方式的材料，生成一份可以分享的ppt"


# ---------------------------------------------------------------------------
# env loader — pull keys from baizhi-agent/.baizhi-agent/config/.env
# ---------------------------------------------------------------------------


def _load_env_file(path: Path) -> None:
    if not path.exists():
        return
    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        os.environ.setdefault(key.strip(), value.strip().strip('"').strip("'"))


# ---------------------------------------------------------------------------
# MCP connect retry — macOS-specific EADDRNOTAVAIL mitigation (test-only)
#
# We observed ~25% of live runs flaking at MCP HTTP transport setup with
# `[Errno 49] Can't assign requested address` wrapped by anyio's TaskGroup.
# It's an OS-level source-port exhaustion glitch, not an MCP issue. LiteLLM's
# `num_retries` doesn't help because MCP isn't routed through LiteLLM.
# We absorb it here with a thin subclass that retries `.connect()`. spec § 7.2
# leaves lifecycle to the use site, so this stays out of the SDK proper.
# ---------------------------------------------------------------------------


class _RetryingMcpToolset(McpToolset):
    """Retries `connect()` on transient EADDRNOTAVAIL / TaskGroup-wrapped variants."""

    async def connect(self) -> None:
        last_exc: BaseException | None = None
        for attempt in range(3):
            try:
                await super().connect()
                return
            except BaseException as exc:  # noqa: BLE001
                if not _is_transient_connect_error(exc):
                    raise
                last_exc = exc
                if attempt < 2:
                    await asyncio.sleep(0.5 * (2 ** attempt))
        assert last_exc is not None
        raise last_exc


def _is_transient_connect_error(exc: BaseException) -> bool:
    if isinstance(exc, BaseExceptionGroup):
        return any(_is_transient_connect_error(sub) for sub in exc.exceptions)
    if isinstance(exc, OSError) and exc.errno in {49, 60, 61, 65, 110}:
        # EADDRNOTAVAIL / ETIMEDOUT / ECONNREFUSED / ENETUNREACH
        return True
    if isinstance(exc, (TimeoutError, asyncio.TimeoutError)):
        return True
    return False


# ---------------------------------------------------------------------------
# Business toolset — write a real .pptx with python-pptx
# ---------------------------------------------------------------------------


class _PptxDeckToolset(BaseToolset):
    name = "pptx_deck_writer"

    def build_schemas(self) -> list[ToolSchema]:
        return [
            ToolSchema(
                name="create_pptx_deck",
                description=(
                    "Create a real .pptx deck file from researched slide content. "
                    "Call this only after loading the pptx skill and collecting "
                    "WebSearch evidence."
                ),
                parameters={
                    "type": "object",
                    "properties": {
                        "filename": {
                            "type": "string",
                            "description": "Output filename ending in .pptx.",
                        },
                        "title": {"type": "string"},
                        "subtitle": {"type": "string"},
                        "slides": {
                            "type": "array",
                            "minItems": 5,
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string"},
                                    "bullets": {
                                        "type": "array",
                                        "items": {"type": "string"},
                                    },
                                    "notes": {"type": "string"},
                                },
                                "required": ["title", "bullets"],
                            },
                        },
                        "sources": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "URLs or source names used in research.",
                        },
                    },
                    "required": ["filename", "title", "slides"],
                    "additionalProperties": False,
                },
            )
        ]

    async def execute(self, call: ToolCall, ctx: ToolCallContext) -> ToolResult:
        args = call.arguments or {}
        filename = str(args.get("filename") or "anthropic-ai-native-organization.pptx")
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        if "/" in filename or "\\" in filename:
            return ToolResult(call_id=call.id, content="ERROR: filename must be a basename", is_error=True)
        slides = args.get("slides") or []
        if not isinstance(slides, list) or len(slides) < 3:
            return ToolResult(call_id=call.id, content="ERROR: provide at least 3 slides", is_error=True)
        output_dir = ctx.storage / "e2e-output"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / filename
        _write_minimal_pptx(
            output_path,
            title=str(args.get("title") or "Anthropic AI Native Organization"),
            subtitle=str(args.get("subtitle") or TASK),
            slides=slides,
            sources=[str(s) for s in (args.get("sources") or [])],
        )
        ctx.run_state["pptx_path"] = str(output_path)
        return ToolResult(
            call_id=call.id,
            content=json.dumps(
                {
                    "pptx_path": str(output_path),
                    "size_bytes": output_path.stat().st_size,
                    "slides": len(slides) + 1,
                },
                ensure_ascii=False,
            ),
        )


def _write_minimal_pptx(
    output_path: Path,
    *,
    title: str,
    subtitle: str,
    slides: list[Any],
    sources: list[str],
) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = title_slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12), Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.LEFT
    title_p.runs[0].font.size = Pt(36)
    title_p.runs[0].font.bold = True
    title_p.runs[0].font.color.rgb = RGBColor(31, 41, 55)
    subtitle_box = title_slide.shapes.add_textbox(Inches(0.75), Inches(2.35), Inches(11.5), Inches(1.4))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].runs[0].font.size = Pt(18)
    subtitle_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(75, 85, 99)

    for raw_slide in slides:
        if not isinstance(raw_slide, dict):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        heading = str(raw_slide.get("title") or "Untitled")[:100]
        bullets = raw_slide.get("bullets") or []
        if not isinstance(bullets, list):
            bullets = [str(bullets)]

        heading_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12.0), Inches(0.65))
        heading_tf = heading_box.text_frame
        heading_tf.text = heading
        heading_run = heading_tf.paragraphs[0].runs[0]
        heading_run.font.size = Pt(28)
        heading_run.font.bold = True
        heading_run.font.color.rgb = RGBColor(17, 24, 39)

        body_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.55), Inches(11.3), Inches(5.3))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.clear()
        for index, bullet in enumerate(bullets[:6]):
            paragraph = body_tf.paragraphs[0] if index == 0 else body_tf.add_paragraph()
            paragraph.text = str(bullet)[:220]
            paragraph.level = 0
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(31, 41, 55)

    if sources:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        heading_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12.0), Inches(0.65))
        heading_tf = heading_box.text_frame
        heading_tf.text = "Sources"
        heading_tf.paragraphs[0].runs[0].font.size = Pt(28)
        heading_tf.paragraphs[0].runs[0].font.bold = True
        body_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.4), Inches(11.5), Inches(5.6))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        body_tf.clear()
        for index, source in enumerate(sources[:8]):
            paragraph = body_tf.paragraphs[0] if index == 0 else body_tf.add_paragraph()
            paragraph.text = str(source)[:180]
            paragraph.font.size = Pt(12)

    prs.save(output_path)


# ---------------------------------------------------------------------------
# The test — note how concise it is post-§17. Hello-world style.
# ---------------------------------------------------------------------------


@pytest.mark.live
@pytest.mark.asyncio
async def test_live_llm_websearch_and_pptx_artifact(tmp_path: Path) -> None:
    if os.environ.get("RUN_BAIZHI_E2E") != "1":
        pytest.skip("set RUN_BAIZHI_E2E=1 to spend real LLM/WebSearch quota")
    _load_env_file(BAIZHI_ENV_FILE)
    if not os.environ.get("MINIMAX_API_KEY"):
        pytest.skip("MINIMAX_API_KEY missing in baizhi-agent/.baizhi-agent/config/.env")
    if not os.environ.get("DASHSCOPE_API_KEY"):
        pytest.skip("DASHSCOPE_API_KEY missing in baizhi-agent/.baizhi-agent/config/.env")

    # MiniMax is OpenAI-compatible — LiteLLM routes via the `openai/` prefix
    # with explicit `api_base`. `num_retries=2` handles LLM-side transient errors
    # (502 / 503 / network blips). MCP-side retry stays in _RetryingMcpToolset.
    model_name = os.environ.get("MINIMAX_MODEL", "MiniMax-M2.7")
    model = LiteLlm(
        f"openai/{model_name}",
        api_base=os.environ.get("MINIMAX_BASE_URL", "https://api.minimaxi.com/v1"),
        api_key=os.environ["MINIMAX_API_KEY"],
        num_retries=2,
    )

    agent = Agent(
        name="baizhi-pptx-deck",
        model=model,
        instruction=(
            "You are running a live E2E test. You MUST use tools before the final answer. "
            "Required sequence: first call load_skill for pptx, then call "
            "load_skill_resource for pptxgenjs.md, then call the WebSearch MCP tool "
            "to research Anthropic AI-native organization / operating model, then "
            "call create_pptx_deck to write the actual .pptx artifact. "
            "Do not claim completion until create_pptx_deck returns a pptx_path. "
            "Keep the deck concise, 5-8 content slides, Chinese language, with sources."
        ),
        tools=[
            SkillCatalogToolset(
                FilesystemSkillRegistry(BUNDLED_SKILLS_ROOT),
                tenant_id="tenant-baizhi-e2e",
            ),
            _RetryingMcpToolset.http(
                "web-search",
                url=WEBSEARCH_URL,
                headers={"Authorization": "Bearer ${DASHSCOPE_API_KEY}"},
                connect_timeout=60,
            ),
            _PptxDeckToolset(),
        ],
        workspace_root=tmp_path / "runs",
        storage_root=tmp_path / "storage",
        default_max_rounds=12,
        default_temperature=0.1,
        default_max_tokens=int(os.environ.get("MINIMAX_MAX_TOKENS", "4096")),
    )

    result = await agent.run(
        TASK,
        tenant_id="tenant-baizhi-e2e",
        enabled_skills=["pptx"],
    )

    tool_names = [
        event.payload["name"]
        for event in result.events
        if event.kind == "tool_call"
    ]
    assert "load_skill" in tool_names, tool_names
    assert "load_skill_resource" in tool_names, tool_names
    assert any(name.startswith("mcp__web-search__") for name in tool_names), tool_names
    assert "create_pptx_deck" in tool_names, tool_names

    outputs = list((tmp_path / "storage" / "e2e-output").glob("*.pptx"))
    assert outputs, f"no pptx output produced; final_text={result.final_text!r}"
    pptx_path = outputs[0]
    assert pptx_path.stat().st_size > 1000
    from pptx import Presentation

    reopened = Presentation(pptx_path)
    assert len(reopened.slides) >= 5
    extracted_text = "\n".join(
        shape.text
        for slide in reopened.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "Anthropic" in extracted_text
